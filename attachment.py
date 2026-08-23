import discord
import os
import mimetypes
from console import console
import base64
import aiohttp
import asyncio
import tempfile
from PIL import Image
from io import BytesIO
import magic
from urllib.parse import urlparse
from utils.file_converter import voicebank_to_json

# Shared session for attachment functions: use central SessionManager
from utils.http_session import session_manager

async def _get_shared_session():
    return await session_manager.get_session()


def clean_url(url: str) -> str:
    """Strip expiring token params from Discord CDN URLs (ex, is, hm)."""
    parsed = urlparse(url)
    return parsed._replace(query="").geturl()


MODEL_EXTS = ('.obj', '.fbx', '.gltf', '.glb', '.dae', '.3ds', '.stl', '.blend')
ZIP_EXTS = ('.zip', '.rar', '.7z', '.tar', '.gz')
MIDI_EXTS = ('.mid', '.midi')
TEXT_EXTENSIONS = (
    ".txt", ".log", ".ini", ".cfg", ".conf", ".properties", ".env", ".md", ".rst", ".adoc",
    ".json", ".csv", ".tsv", ".xml", ".yaml", ".yml", ".toml", ".xsd", ".dtd", 
    ".html", ".htm", ".xhtml", ".css",
    ".py", ".ipynb", ".js", ".mjs", ".cjs", ".ts", ".tsx", ".jsx", ".java", ".kt", ".kts",
    ".c", ".h", ".cpp", ".hpp", ".cc", ".cxx", ".hh", ".hxx", ".cs", ".go", ".rs", ".php",
    ".phtml", ".php3", ".php4", ".php5", ".phps",
    ".rb", ".swift", ".scala", ".sc", ".dart", ".pl", ".pm", ".t", ".r", ".rmd", ".jl",
    ".sh", ".bash", ".zsh", ".ksh", ".fish", ".bat", ".cmd", ".ps1", ".psm1", ".psd1",
    ".lua", ".sql", ".erl", ".hrl", ".ex", ".exs", ".clj", ".cljs", ".edn",
    ".hs", ".lhs", ".ml", ".mli", ".mll", ".mly", ".vb", ".vbs", ".f", ".for", ".f90",
    ".f95", ".f03", ".f08", ".asm", ".s", ".v", ".vhd", ".vhdl", ".sv", ".list", ".mk"
)
AUDIO_EXTS = ('.mp3', '.wav', '.m4a', '.m4b', '.m4r', '.flac', '.ogg', '.aac')
VIDEO_EXTS = ('.mp4', '.mkv', '.avi', '.mov', '.wmv', '.flv', '.webm')
IMG_EXTS = ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.svg', '.webp')

lang_to_ext = {
    "python": "py",
    "javascript": "js",
    "typescript": "ts",
    "java": "java",
    "c": "c",
    "c++": "cpp",
    "c#": "cs",
    "go": "go",
    "rust": "rs",
    "kotlin": "kt",
    "swift": "swift",
    "ruby": "rb",
    "php": "php",
    "html": "html",
    "css": "css",
    "sql": "sql",
    "shell": "sh",
    "bash": "sh",
    "powershell": "ps1",
    "r": "r",
    "lua": "lua",
    "perl": "pl",
    "scala": "scala",
    "dart": "dart",
    "objective-c": "m",
    "objective-c++": "mm",
    "haskell": "hs",
    "elixir": "ex",
    "erlang": "erl",
    "json": "json",
    "yaml": "yml",
    "toml": "toml",
    "ini": "ini",
    "xml": "xml",
    "markdown": "md",
    "plaintext": "txt",
    "assembly": "asm",
    "verilog": "v",
    "vhdl": "vhd",
    "matlab": "m",
    "julia": "jl",
    "fortran": "f90",
    "vb.net": "vb",
    "groovy": "groovy",
    "nim": "nim",
    "zig": "zig",
    "solidity": "sol",
    "dockerfile": "dockerfile",
    "makefile": "mk",
    "protobuf": "proto",
    "graphql": "graphql",
    "latex": "tex",
}

def is_video_attachment(att: discord.Attachment) -> bool:
  if att.content_type and att.content_type.startswith("video/"):
    return True
  video_extensions = VIDEO_EXTS
  ext = os.path.splitext(att.filename)[1].lower()
  if ext in video_extensions:
    return True
  if att.size > 10 * 1024 * 1024:
    return True
  return False

PDF_EXTENSIONS = [".pdf"]

def is_pdf_attachment(att: discord.Attachment) -> bool:    
  mime_type = att.content_type or mimetypes.guess_type(att.filename)[0] or ""
  ext = os.path.splitext(att.filename)[1].lower()  
  if "application/pdf" in mime_type:
    return True
  if ext in PDF_EXTENSIONS:
    return True
      
  return False

def is_image_attachment(att: discord.Attachment) -> bool:
  ct = (att.content_type or "").lower()
  name = (att.filename or "").lower()
  return ("image" in ct) or name.endswith(IMG_EXTS)

async def is_text_attachment(att: discord.Attachment) -> bool:
  """Detect whether an attachment is text.

  Strategy:
  - Fast checks: filename extension, content-type hints (json/xml/yaml or text/*)
  - If ambiguous, fetch a small byte range (0-4095) and try to decode as UTF-8 (strict). If decoding succeeds, treat as text.
  """
  name = (att.filename or "").lower()
  ct = (att.content_type or "").lower()

  # Fast heuristics
  if (
    name.endswith(TEXT_EXTENSIONS)
    or ct.startswith("text/")
    or "json" in ct
    or "xml" in ct
    or "yaml" in ct or "yml" in ct
  ):
    return True

  # If size is very large, avoid fetching; assume non-text
  try:
    if getattr(att, 'size', None) and att.size > 5 * 1024 * 1024:
      return False
  except Exception:
    pass

  # Fallback: fetch small prefix and test UTF-8 validity
  try:
    headers = {"Range": "bytes=0-4095", "User-Agent": "Arona/1.0"}
    timeout = aiohttp.ClientTimeout(total=5)
    session = await _get_shared_session()
    async with session.get(att.url, headers=headers, timeout=timeout) as resp:
      if resp.status not in (200, 206):
        return False
      data = await resp.content.read(4096)
      if not data:
        return False
      try:
        data.decode('utf-8')
        return True
      except UnicodeDecodeError:
        return False
  except Exception as e:
    console.log(f"is_text_attachment encoding detect failed for {att.filename}: {e}", "DEBUG")
    return False
  
def is_audio_attachment(att: discord.Attachment) -> bool:
  name = (att.filename or "").lower()
  ct = (att.content_type or "").lower()
  return name.endswith(AUDIO_EXTS) or "audio" in ct

def is_model3d_attachment(att: discord.Attachment) -> bool:
  name = (att.filename or "").lower()
  ct = (att.content_type or "").lower()
  return name.endswith(MODEL_EXTS) or "model" in ct or "3d" in ct

def is_zip_attachment(att: discord.Attachment) -> bool:
  name = (att.filename or "").lower() 
  ct = (att.content_type or "").lower()
  return name.endswith(ZIP_EXTS) or "zip" in ct or "archive" in ct or "compressed" in ct

def is_midi_attachment(att: discord.Attachment) -> bool:
  name = (att.filename or "").lower()
  ct = (att.content_type or "").lower()
  return name.endswith(MIDI_EXTS) or "midi" in ct

def is_voicebank_attachment(att: discord.Attachment) -> bool:
  name = (att.filename or "").lower()
  return name.endswith(('.ust', '.vsqx'))

def is_weight_attachment(att: discord.Attachment) -> bool:
  name = (att.filename or "").lower()
  ct = (att.content_type or "").lower()
  return name.endswith(('.pth', '.pt', '.bin')) or "model" in ct or "weight" in ct

def is_numpy_attachment(att: discord.Attachment) -> bool:
  name = (att.filename or "").lower()
  ct = (att.content_type or "").lower()
  return name.endswith(('.npy', '.npz')) or "numpy" in ct

GEMINI_SUPPORTED_MIME_TYPES = {
    # Images
    "image/jpeg",
    "image/png",
    "image/webp",
    
    # Video
    "video/mp4",
    "video/mpeg",
    "video/quicktime",
    "video/webm",
    "video/x-msvideo",
    
    # Audio
    "audio/aac",
    "audio/flac",
    "audio/mp3",
    "audio/mpeg",
    "audio/opus",
    "audio/ogg",
    "audio/wav",
    "audio/webm",
    
    # Text
    "text/plain",
    
    # PDF
    "application/pdf",
}

def get_mime_type_from_magic(file_data: bytes, filename: str = "") -> str:
  """
  Detect MIME type using python-magic library.
  Falls back to filename extension if magic fails.
  """
  try:
    detected_mime = magic.from_buffer(file_data, mime=True)
    if detected_mime:
      return detected_mime.lower().split(';')[0].strip()
  except Exception as e:
    console.log(f"Magic library failed to detect MIME type: {e}", "WARN")
  
  if filename:
    ext = os.path.splitext(filename)[1].lower()
    guess_mime = mimetypes.guess_type(filename)[0]
    if guess_mime:
      return guess_mime.lower()
  
  return "application/octet-stream"

def normalize_mime_type(mime_type: str, filename: str = "") -> str | None:
  """
  Normalize MIME type to Gemini-supported format.
  Returns None if MIME type is not supported.
  """
  mime_type = mime_type.lower().split(';')[0].strip()
  
  if mime_type in GEMINI_SUPPORTED_MIME_TYPES:
    return mime_type
  # Image normalization
  if mime_type.startswith("image/"):
    if "jpeg" in mime_type or "jpg" in mime_type:
      return "image/jpeg"
    if "png" in mime_type:
      return "image/png"
    if "gif" in mime_type:
      return "image/png"  # Convert GIF to PNG (GIF not supported by Gemini)
    if "webp" in mime_type:
      return "image/webp"
  
  # Video normalization
  if mime_type.startswith("video/"):
    if "mp4" in mime_type or "quicktime" in mime_type:
      return "video/mp4"
    if "mpeg" in mime_type:
      return "video/mpeg"
    if "webm" in mime_type:
      return "video/webm"
    if "x-msvideo" in mime_type or "avi" in mime_type:
      return "video/x-msvideo"
  
  # Audio normalization
  if mime_type.startswith("audio/"):
    if "mpeg" in mime_type or "mp3" in mime_type:
      return "audio/mpeg"
    if "wav" in mime_type:
      return "audio/wav"
    if "ogg" in mime_type:
      return "audio/ogg"
    if "flac" in mime_type:
      return "audio/flac"
    if "aac" in mime_type:
      return "audio/aac"
    if "webm" in mime_type:
      return "audio/webm"
    if "opus" in mime_type or "x-opus" in mime_type:
      return "audio/opus"
  
  # Text normalization
  if (
      mime_type.startswith("text/")
      or any(k in mime_type for k in ("json", "xml", "yaml", "yml"))
      or (filename and filename.lower().endswith(TEXT_EXTENSIONS))
  ):
    return "text/plain"
  
  # PDF
  if "pdf" in mime_type:
    return "application/pdf"
  
  return None

# Semaphore để giới hạn số lượng xử lý đồng thời
CONCURRENT_LIMIT = 5
attachment_semaphore = asyncio.Semaphore(CONCURRENT_LIMIT)

def _unsupported_attachment_part(att: discord.Attachment, reason: str) -> list[dict]:
    name = (att.filename or "").lower()
    ext = os.path.splitext(name)[1]
    url = att.url
    local_name = att.filename  # keep original name for curl

    # per-type download + parse hints
    if ext in ('.xlsx', '.xls', '.xlsm', '.ods'):
        tool_hint = "run_code (pandas + openpyxl)"
        read_code = (
            "import pandas as pd\n"
            f"sheets = pd.read_excel('{local_name}', sheet_name=None)  # None = all sheets\n"
            "for sheet_name, df in sheets.items():\n"
            "    print(f'=== {{sheet_name}} ===')\n"
            "    print(df.head())\n"
            "    print(df.dtypes)\n"
            "    print()"
        )
    elif ext == '.csv':
        tool_hint = "run_code (pandas)"
        read_code = (
            "import pandas as pd\n"
            f"df = pd.read_csv('{local_name}')\n"
            "print(df.head())\n"
            "print(df.dtypes)"
        )
    elif ext in ('.docx', '.doc'):
        tool_hint = "run_code (python-docx)"
        read_code = (
            "from docx import Document\n"
            f"doc = Document('{local_name}')\n"
            "for p in doc.paragraphs:\n"
            "    if p.text.strip():\n"
            "        print(p.text)"
        )
    elif ext in ('.pptx', '.ppt'):
        tool_hint = "run_code (python-pptx)"
        read_code = (
            "from pptx import Presentation\n"
            f"prs = Presentation('{local_name}')\n"
            "for i, slide in enumerate(prs.slides):\n"
            "    print(f'--- Slide {{i+1}} ---')\n"
            "    for shape in slide.shapes:\n"
            "        if shape.has_text_frame:\n"
            "            print(shape.text)"
        )
    elif ext == '.zip':
        tool_hint = "run_shell (unzip) or run_code (zipfile)"
        read_code = (
            "import zipfile, os\n"
            f"with zipfile.ZipFile('{local_name}', 'r') as z:\n"
            "    print('Contents:', z.namelist())\n"
            "    z.extractall('extracted/')\n"
            "# Files now in extracted/"
        )
    elif ext in ('.tar', '.gz', '.bz2', '.xz', '.tgz'):
        tool_hint = "run_shell (tar)"
        read_code = f"# run_shell: tar -xf '{local_name}' -C extracted/"
    elif ext == '.rar':
        tool_hint = "run_shell (unrar)"
        read_code = f"# run_shell: unrar x '{local_name}' extracted/"
    elif ext == '.7z':
        tool_hint = "run_shell (7z)"
        read_code = f"# run_shell: 7z x '{local_name}' -o extracted/"
    elif ext in ('.sqlite', '.db', '.sqlite3'):
        tool_hint = "run_code (sqlite3)"
        read_code = (
            "import sqlite3\n"
            f"conn = sqlite3.connect('{local_name}')\n"
            "tables = conn.execute(\"SELECT name FROM sqlite_master WHERE type='table'\").fetchall()\n"
            "print('Tables:', [t[0] for t in tables])\n"
            "for t in tables:\n"
            "    print(f'--- {{t[0]}} ---')\n"
            "    print(conn.execute(f'SELECT * FROM {{t[0]}} LIMIT 5').fetchall())"
        )
    elif ext == '.parquet':
        tool_hint = "run_code (pandas)"
        read_code = (
            "import pandas as pd\n"
            f"df = pd.read_parquet('{local_name}')\n"
            "print(df.head())\n"
            "print(df.dtypes)"
        )
    elif ext in ('.pkl', '.pickle'):
        tool_hint = "run_code (pickle)"
        read_code = (
            "import pickle\n"
            f"with open('{local_name}', 'rb') as f:\n"
            "    data = pickle.load(f)\n"
            "print(type(data))\n"
            "print(data)"
        )
    elif ext == '.json':
        tool_hint = "run_code (json)"
        read_code = (
            "import json\n"
            f"with open('{local_name}', 'r', encoding='utf-8') as f:\n"
            "    data = json.load(f)\n"
            "print(type(data))\n"
            "import pprint; pprint.pprint(data)"
        )
    else:
        tool_hint = "run_shell or run_code"
        read_code = (
            "# Inspect file type first:\n"
            "import subprocess\n"
            f"print(subprocess.run(['file', '{local_name}'], capture_output=True, text=True).stdout)"
        )

    hint_text = (
        f"[Attachment cannot be read directly: {att.filename}]\n"
        f"Reason: {reason}\n"
        f"\n"
        f"To access this file, use {tool_hint}:\n"
        f"\n"
        f"Step 1 — download via run_shell:\n"
        f"  curl -L -o '{local_name}' '{url}'\n"
        f"\n"
        f"Step 2 — read/process via run_code:\n"
        f"{read_code}"
    )
    return [
        {"text": f"[Attachment: {att.filename} | URL: {att.url}]"},
        {
            "inline_data": {
                "mime_type": "text/plain",
                "data": base64.b64encode(hint_text.encode('utf-8')).decode('utf-8')
            }
        }
    ]


async def process_single_attachment(att: discord.Attachment, text: bool = True) -> list[dict] | None:
    """
    Xử lý 1 attachment độc lập
    Returns: list[dict] part hoặc None nếu thất bại
    """
    async with attachment_semaphore:  # Giới hạn concurrent processing
        try:
            # 0. Pre-check size trước khi download (att.size có sẵn từ Discord metadata,
            #    khỏi tốn băng thông + thời gian tải file chắc chắn sẽ bị reject)
            _MAX_POSSIBLE_LIMIT = 256 * 1024 * 1024  # limit lớn nhất trong _SIZE_LIMITS (PDF)
            if getattr(att, 'size', None) and att.size > _MAX_POSSIBLE_LIMIT:
                console.log(
                    f"Skip download, file too large: {att.filename} "
                    f"({att.size/1024/1024:.1f}MB)", "WARN"
                )
                return _unsupported_attachment_part(
                    att, f"file too large ({att.size/1024/1024:.1f}MB) to process"
                )

            # 1. Download với retry
            max_retries = 3
            raw_data = None
            
            for attempt in range(max_retries):
                try:
                    raw_data = await asyncio.wait_for(att.read(), timeout=30.0)
                    break
                except asyncio.TimeoutError:
                    if attempt == max_retries - 1:
                        console.log(f"Timeout downloading {att.filename}", "WARN")
                        return _unsupported_attachment_part(att, "timeout downloading")
                    await asyncio.sleep(1)
                except Exception as e:
                    if attempt == max_retries - 1:
                        console.log(f"Failed to download {att.filename}: {e}", "WARN")
                        return _unsupported_attachment_part(att, f"download error: {e}")
                    await asyncio.sleep(1)
            
            if not raw_data:
                return _unsupported_attachment_part(att, "empty download")
            
            # 1b. Block known binary extensions that masquerade as supported MIME types
            BLACKLISTED_EXTS = ('.ai', '.eps', '.ps')
            if att.filename.lower().endswith(BLACKLISTED_EXTS):
                ext = os.path.splitext(att.filename)[1]
                console.log(f"Blocked unsupported file type: {att.filename} ({ext})", "WARN")
                return _unsupported_attachment_part(att, f"unsupported file type: {ext}")
            
            # 2. Detect MIME type (offload sang thread — magic.from_buffer là blocking C call,
            #    chạy trực tiếp trong event loop sẽ đóng băng cả bot trong lúc detect)
            mime_type = await asyncio.to_thread(get_mime_type_from_magic, raw_data, att.filename)

            # 2b. Audio-misidentified-as-video intercept
            #     Một số format audio dùng container video (M4A/M4B/M4R = MPEG-4,
            #     OGG audio đôi khi bị detect là video/ogg, WebM audio → video/webm, v.v.)
            #     python-magic và Discord đều bị lừa → báo video/* → Gemini reject
            #     "0 frames found". Convert sang WAV bằng ffmpeg trước khi normalize.
            _fname_lower = att.filename.lower()
            _src_ext = os.path.splitext(_fname_lower)[1]
            if _fname_lower.endswith(AUDIO_EXTS) and mime_type.startswith('video/'):
                try:
                    from utils.file_converter import audio_to_wav
                    console.log(
                        f"Audio ext ({_src_ext}) misdetected as {mime_type}, "
                        f"converting to WAV: {att.filename}", "INFO"
                    )
                    raw_data = await audio_to_wav(raw_data, src_ext=_src_ext)
                    mime_type = 'audio/wav'
                except Exception as _conv_err:
                    console.log(f"audio→WAV failed for {att.filename}: {_conv_err}", "WARN")
                    return _unsupported_attachment_part(att, f"audio conversion failed: {_conv_err}")

            # 3. Skip text nếu không cần -> nhưng vẫn thông báo attachment
            if mime_type == "text/plain" and not text:
                console.log(f"Skipping text attachment (not requested): {att.filename}", "INFO")
                return _unsupported_attachment_part(att, "text attachment omitted by request")
            
            # 4. Normalize MIME
            normalized_mime = normalize_mime_type(mime_type, att.filename)
            if not normalized_mime:
                console.log(f"Unsupported MIME: {att.filename} ({mime_type})", "WARN")
                return _unsupported_attachment_part(att, f"unsupported MIME: {mime_type}")
              
            if normalized_mime == "text/plain" and not text:
                console.log(f"Skipping text attachment (not requested): {att.filename}", "INFO")
                return _unsupported_attachment_part(att, "text attachment omitted by request")
            
            # 4b. Size limit check per MIME category
            _SIZE_LIMITS = {
                "application/pdf": 256 * 1024 * 1024,
                "video/":          128 * 1024 * 1024,
                "audio/":          80 * 1024 * 1024,
                "image/":          40 * 1024 * 1024,
                "text/plain":       30 * 1024 * 1024,
            }
            _file_size = len(raw_data)
            for _prefix, _limit in _SIZE_LIMITS.items():
                if normalized_mime == _prefix or normalized_mime.startswith(_prefix):
                    if _file_size > _limit:
                        console.log(
                            f"File too large for inline_data: {att.filename} "
                            f"({_file_size/1024/1024:.1f}MB > {_limit/1024/1024:.0f}MB limit for {normalized_mime})", "WARN"
                        )
                        return _unsupported_attachment_part(
                            att, f"file too large ({_file_size/1024/1024:.1f}MB, limit {_limit/1024/1024:.0f}MB for {normalized_mime})"
                        )
                    break

            # 5. Convert nếu cần
            final_mime = normalized_mime
            final_data = raw_data
            
            # GIF -> PNG
            if att.filename.lower().endswith('.gif'):
                from utils.file_converter import gif_to_png
                final_data = await gif_to_png(raw_data)
                final_mime = "image/png"
                console.log(f"Converted GIF: {att.filename}", "INFO")
            
            # MIDI -> WAV
            elif mime_type in ("audio/midi", "audio/x-midi"):
                from utils.file_converter import midi_to_wav
                final_data = await midi_to_wav(raw_data)
                final_mime = "audio/wav"
                console.log(f"Converted MIDI: {att.filename}", "INFO")
            
            # 3D Model -> Image
            elif att.filename.lower().endswith(('.obj', '.fbx', '.gltf', '.glb', '.stl')):
                from utils.file_converter import model3d_to_image
                final_data = await model3d_to_image(raw_data)
                final_mime = "image/png"
                console.log(f"Converted 3D model: {att.filename}", "INFO")
            
            # Voicebank -> JSON
            elif att.filename.lower().endswith(('.ust', '.vsqx')):
                from utils.file_converter import voicebank_to_json
                file_type = "ust" if att.filename.lower().endswith('.ust') else "vsqx"
                json_data = await voicebank_to_json(raw_data, file_type)
                final_data = json_data.encode('utf-8')
                final_mime = "text/plain"
                console.log(f"Converted voicebank: {att.filename}", "INFO")
            
            # Weight files -> Text
            elif att.filename.lower().endswith(('.pth', '.pt', '.bin')):
                from utils.file_converter import convert_model_to_text_bytes
                final_data = await convert_model_to_text_bytes(raw_data)
                final_mime = "text/plain"
                console.log(f"Converted model: {att.filename}", "INFO")
            
            # NumPy -> Text
            elif att.filename.lower().endswith(('.npy', '.npz')):
                from utils.file_converter import numpy_to_text
                final_data = await numpy_to_text(raw_data, att.filename)
                final_mime = "text/plain"
                console.log(f"Converted NumPy: {att.filename}", "INFO")
            
            # 6. Wrap text files in code block trước khi encode
            if final_mime == "text/plain":
                ext = os.path.splitext(att.filename)[1].lower().lstrip('.')
                ext_to_lang = {v: k for k, v in lang_to_ext.items()}
                lang = ext_to_lang.get(ext, ext or "")
                text_content = final_data.decode('utf-8', errors='replace')

                # truncate at 2000 lines and stage full file for read_file
                TEXT_LINE_LIMIT = 2000
                all_lines = text_content.splitlines(keepends=True)
                total_lines = len(all_lines)
                truncated = False
                file_id_hint = None

                if total_lines > TEXT_LINE_LIMIT:
                    truncated = True
                    # Stage the full file so model can call read_file
                    try:
                        from utils.edit_text_file import create_files as _stage_files
                        staged = _stage_files([{"filename": att.filename, "content": text_content}])
                        if staged and staged[0].get("file_id"):
                            file_id_hint = staged[0]["file_id"]
                    except Exception as _e:
                        console.log(f"[attachment] Failed to stage full file {att.filename}: {_e}", "WARN")
                    # Trim to limit
                    text_content = "".join(all_lines[:TEXT_LINE_LIMIT])

                wrapped = f"```{lang}\n{text_content}\n```"
                if truncated:
                    hint = (
                        f"\n[FILE TRUNCATED — showing lines 1–{TEXT_LINE_LIMIT} of {total_lines}. "
                        f"Full file staged as file_id: `{file_id_hint}`. "
                        f"Call read_file(file_ref=\"{file_id_hint}\", start_line=..., end_line=...) to read more.]"
                    )
                    wrapped += hint
                final_data = wrapped.encode('utf-8')

            # 7. Encode to base64 (offload sang thread nếu file đủ lớn — encode đồng bộ
            #    file video/audio hàng chục MB có thể block event loop cả giây)
            if len(final_data) > 2 * 1024 * 1024:
                base64_data = await asyncio.to_thread(
                    lambda d: base64.b64encode(d).decode("utf-8"), final_data
                )
            else:
                base64_data = base64.b64encode(final_data).decode("utf-8")
            
            # 8. Log preview for images
            if "image" in final_mime:
                console.log(f'Att preview: <img src="{att.url}" width="70" height="auto">', "DEBUG")
            
            console.log(f"Loaded: {att.filename} ({final_mime}), {len(final_data)} bytes", "INFO")
            
            return [
                {"text": f"[Attachment: {att.filename} | URL: {att.url}]"},
                {
                    "inline_data": {
                        "mime_type": final_mime,
                        "data": base64_data
                    }
                }
            ]
            
        except Exception as e:
            console.log(f"Error processing {att.filename}: {e}", "ERROR")
            return _unsupported_attachment_part(att, f"processing error: {e}")


async def process_archive_attachment(att: discord.Attachment) -> list[dict]:
    """
    Archive không được giải nén tự động.
    Trả fallback alert model tự tải file với curl và xử lý tay.
    """
    async with attachment_semaphore:
        return _unsupported_attachment_part(att, "archive extraction disabled; use curl + run_code to inspect")


async def discord_attachment_to_parts(attachments: list[discord.Attachment], text=True) -> list[dict]:
    """
    Convert Discord attachments to Gemini-compatible format với parallel processing.
    
    Improvements:
    - Xử lý song song tất cả attachments (giới hạn bởi semaphore)
    - Archive files được xử lý riêng vì tạo nhiều parts
    - Tốc độ tăng 3-5x với nhiều attachments
    """
    if not attachments:
        return []
    
    # Phân loại attachments
    archive_atts = []
    normal_atts = []
    
    for att in attachments:
        name = att.filename.lower()
        mime = att.content_type or ""
        
        is_archive = (
            name.endswith(('.zip', '.rar', '.7z', '.tar', '.gz', '.bz2', '.xz')) or
            "zip" in mime or "archive" in mime or "compressed" in mime
        )
        
        if is_archive:
            archive_atts.append(att)
        else:
            normal_atts.append(att)
    
    # Xử lý song song
    tasks = []
    
    # Normal attachments - mỗi att = 1 part
    tasks.extend([process_single_attachment(att, text) for att in normal_atts])
    
    # Archive attachments - mỗi att = nhiều parts
    tasks.extend([process_archive_attachment(att) for att in archive_atts])
    
    console.log(f"Processing {len(attachments)} attachments...", "INFO")
    start_time = asyncio.get_event_loop().time()
    
    results = await asyncio.gather(*tasks, return_exceptions=True)
        
    # Flatten results
    parts = []
    for result in results:
        if isinstance(result, Exception):
            console.log(f"Task exception: {result}", "ERROR")
            continue
        
        if result is None:
            continue
        
        # Archive returns list, single attachment returns dict
        if isinstance(result, list):
            parts.extend(result)
        else:
            parts.append(result)
    
    console.log(f"Total parts: {len(parts)}", "INFO")
    return parts